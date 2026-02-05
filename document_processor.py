"""
Document processing with tree-based indexing.
Self-contained version using OpenAI directly (no external pageindex dependency).
Optimized for large document collections (4000+).
"""
import os
import io
import json
import asyncio
import base64
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path

# Document extraction libraries
import pymupdf  # PyMuPDF (new API)
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import openai

from config import get_settings

settings = get_settings()

# Directory for storing page images
PAGE_IMAGES_DIR = Path("uploads/page_images")
PAGE_IMAGES_DIR.mkdir(parents=True, exist_ok=True)


# ============ PAGE IMAGE EXTRACTION ============

def extract_page_images_from_pdf(file_path: str, doc_id: int, max_pages: int = 50) -> List[str]:
    """Extract pages as images from PDF, return list of image paths.
    
    Args:
        file_path: Path to PDF file
        doc_id: Document ID for organizing images
        max_pages: Maximum number of pages to extract (default 50 for performance)
    """
    doc = pymupdf.open(file_path)
    image_paths = []
    
    doc_images_dir = PAGE_IMAGES_DIR / str(doc_id)
    doc_images_dir.mkdir(parents=True, exist_ok=True)
    
    total_pages = min(len(doc), max_pages)
    
    for page_num, page in enumerate(doc, 1):
        if page_num > max_pages:
            break
        try:
            # Use 1.5x zoom for balance of quality and speed
            mat = pymupdf.Matrix(1.5, 1.5)
            pix = page.get_pixmap(matrix=mat)
            img_path = doc_images_dir / f"page_{page_num}.png"
            pix.save(str(img_path))
            image_paths.append(str(img_path))
        except Exception as e:
            print(f"  Warning: Could not extract page {page_num}: {e}")
            continue
    
    doc.close()
    return image_paths


def extract_page_images_from_pptx(file_path: str, doc_id: int) -> List[str]:
    """Extract slides as images from PowerPoint. Uses text rendering as fallback."""
    from PIL import Image, ImageDraw, ImageFont
    
    doc_images_dir = PAGE_IMAGES_DIR / str(doc_id)
    doc_images_dir.mkdir(parents=True, exist_ok=True)
    
    prs = Presentation(file_path)
    image_paths = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        # Create a simple text-based representation of the slide
        img = Image.new('RGB', (800, 600), color='white')
        draw = ImageDraw.Draw(img)
        
        y_pos = 20
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()[:100]  # Limit text length
                draw.text((20, y_pos), text, fill='black')
                y_pos += 30
                if y_pos > 560:
                    break
        
        img_path = doc_images_dir / f"page_{slide_num}.png"
        img.save(str(img_path))
        image_paths.append(str(img_path))
    
    return image_paths


# ============ TEXT EXTRACTION ============

def extract_text_from_pdf(file_path: str) -> Tuple[str, int, List[Tuple[int, str]]]:
    """Extract text from PDF, return (text, page_count, page_texts)."""
    doc = pymupdf.open(file_path)
    pages = []
    page_texts = []
    for page_num, page in enumerate(doc, 1):
        text = page.get_text()
        pages.append(text)
        page_texts.append((page_num, text))
    doc.close()
    return "\n\n---PAGE BREAK---\n\n".join(pages), len(pages), page_texts


def extract_text_from_docx(file_path: str) -> Tuple[str, int, List[Tuple[int, str]]]:
    """Extract text from Word document as markdown-style text."""
    doc = DocxDocument(file_path)
    lines = []
    page_texts = []
    current_page_text = []
    page_num = 1
    para_count = 0
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name.lower() if para.style else ""
        if "heading 1" in style:
            lines.append(f"# {text}")
            current_page_text.append(f"# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
            current_page_text.append(f"## {text}")
        elif "heading 3" in style:
            lines.append(f"### {text}")
            current_page_text.append(f"### {text}")
        elif "title" in style:
            lines.append(f"# {text}")
            current_page_text.append(f"# {text}")
        else:
            lines.append(text)
            current_page_text.append(text)
        lines.append("")
        para_count += 1
        
        # Approximate page breaks (every ~20 paragraphs)
        if para_count >= 20:
            page_texts.append((page_num, "\n".join(current_page_text)))
            page_num += 1
            current_page_text = []
            para_count = 0
    
    # Add remaining content
    if current_page_text:
        page_texts.append((page_num, "\n".join(current_page_text)))
    
    return "\n".join(lines), page_num, page_texts


def extract_text_from_xlsx(file_path: str) -> Tuple[str, int, List[Tuple[int, str]]]:
    """Extract text from Excel spreadsheet as markdown."""
    wb = load_workbook(file_path, data_only=True)
    lines = []
    page_texts = []
    
    for sheet_idx, sheet_name in enumerate(wb.sheetnames, 1):
        ws = wb[sheet_name]
        sheet_lines = [f"# Sheet: {sheet_name}\n"]
        rows = list(ws.iter_rows(values_only=True))
        for i, row in enumerate(rows):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(v.strip() for v in row_values):
                if i == 0:
                    sheet_lines.append("| " + " | ".join(row_values) + " |")
                    sheet_lines.append("|" + "|".join(["---"] * len(row_values)) + "|")
                else:
                    sheet_lines.append("| " + " | ".join(row_values) + " |")
        lines.extend(sheet_lines)
        lines.append("")
        page_texts.append((sheet_idx, "\n".join(sheet_lines)))
    
    return "\n".join(lines), len(wb.sheetnames), page_texts


def extract_text_from_pptx(file_path: str) -> Tuple[str, int, List[Tuple[int, str]]]:
    """Extract text from PowerPoint as markdown."""
    prs = Presentation(file_path)
    lines = []
    page_texts = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = [f"# Slide {i}\n"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_lines.append(shape.text.strip())
                slide_lines.append("")
        lines.extend(slide_lines)
        lines.append("")
        page_texts.append((i, "\n".join(slide_lines)))
    
    return "\n".join(lines), len(prs.slides), page_texts


def extract_text(file_path: str, file_type: str) -> Tuple[str, int, List[Tuple[int, str]]]:
    """Extract text from document based on type. Returns (text, page_count, page_texts)."""
    extractors = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "xlsx": extract_text_from_xlsx,
        "pptx": extract_text_from_pptx,
    }
    extractor = extractors.get(file_type)
    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type}")
    return extractor(file_path)


def extract_page_images(file_path: str, file_type: str, doc_id: int) -> List[str]:
    """Extract page images from document. Returns list of image paths."""
    if file_type == "pdf":
        return extract_page_images_from_pdf(file_path, doc_id)
    elif file_type == "pptx":
        return extract_page_images_from_pptx(file_path, doc_id)
    # For docx and xlsx, we don't have native image extraction
    return []


# ============ TREE INDEXING ============

async def generate_tree_structure(
    text: str, 
    filename: str, 
    page_count: int,
    page_texts: List[Tuple[int, str]] = None
) -> Dict[str, Any]:
    """
    Generate a hierarchical tree structure from document text.
    Uses OpenAI to analyze structure and create summaries.
    """
    client = openai.AsyncOpenAI(api_key=settings.openai_api_key)
    
    # Truncate text if too long (keep first 100k chars for structure analysis)
    analysis_text = text[:100000] if len(text) > 100000 else text
    
    prompt = f"""Analyze this document and create a hierarchical tree structure.

Document: {filename}
Total pages/slides: {page_count}

For each section, provide:
- node_id: Unique identifier (format: "0001", "0002", etc.)
- title: Section title or topic
- summary: 2-3 sentence summary of the section content
- text: Key content from that section (up to 500 words)
- page_numbers: Array of page/slide numbers where this content appears (e.g., [1, 2] for pages 1 and 2)

Return a JSON structure:
{{
    "doc_name": "{filename}",
    "doc_description": "One paragraph describing what this document is about",
    "structure": [
        {{
            "node_id": "0001",
            "title": "Section Title",
            "summary": "Brief summary...",
            "text": "Key content...",
            "page_numbers": [1],
            "nodes": [
                // Nested subsections if applicable
            ]
        }}
    ]
}}

Document content:
{analysis_text}

Return ONLY valid JSON, no markdown formatting or explanation."""

    try:
        response = await client.chat.completions.create(
            model=settings.openai_model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0,
            response_format={"type": "json_object"}
        )
        
        result = json.loads(response.choices[0].message.content)
        
        # Ensure required fields exist
        if "structure" not in result:
            result["structure"] = []
        if "doc_name" not in result:
            result["doc_name"] = filename
            
        return result
        
    except Exception as e:
        # Fallback: create minimal structure
        print(f"Error generating tree structure: {e}")
        return {
            "doc_name": filename,
            "doc_description": f"Document: {filename}",
            "structure": [{
                "node_id": "0001",
                "title": filename,
                "summary": "Full document content",
                "text": text[:5000],
                "nodes": []
            }]
        }


async def process_document(
    file_path: str, 
    filename: str, 
    file_type: str,
    doc_id: int = None,
    progress_callback = None
) -> Dict[str, Any]:
    """
    Main entry point for document processing.
    Extracts text and generates tree structure.
    progress_callback: Optional async function that takes an int (0-100) for progress updates.
    """
    print(f"Processing document: {filename}")
    
    if progress_callback:
        await progress_callback(10)
    
    # Extract text
    text, page_count, page_texts = extract_text(file_path, file_type)
    print(f"  Extracted {len(text)} chars, ~{page_count} pages")
    
    if progress_callback:
        await progress_callback(30)
    
    # Extract page images if doc_id is provided
    page_image_paths = []
    if doc_id:
        try:
            page_image_paths = extract_page_images(file_path, file_type, doc_id)
            print(f"  Extracted {len(page_image_paths)} page images")
        except Exception as e:
            print(f"  Warning: Could not extract page images: {e}")
    
    if progress_callback:
        await progress_callback(50)
    
    # Generate tree structure with page info
    tree = await generate_tree_structure(text, filename, page_count, page_texts)
    
    if progress_callback:
        await progress_callback(80)
    
    # Add metadata
    tree["_source_file"] = filename
    tree["_file_type"] = file_type
    tree["_char_count"] = len(text)
    tree["_page_count"] = page_count
    tree["_has_page_images"] = len(page_image_paths) > 0
    
    print(f"  Generated tree with {len(tree.get('structure', []))} top-level nodes")
    
    return tree


# ============ UTILITY FUNCTIONS ============

def create_node_mapping(tree: Dict[str, Any], mapping: Dict = None) -> Dict[str, Any]:
    """Create a flat mapping of node_id -> node for quick lookup."""
    if mapping is None:
        mapping = {}
    
    if "structure" in tree:
        nodes = tree["structure"]
        if isinstance(nodes, list):
            for node in nodes:
                _map_node(node, mapping)
        elif isinstance(nodes, dict):
            _map_node(nodes, mapping)
        return mapping
    
    _map_node(tree, mapping)
    return mapping


def _map_node(node: Dict[str, Any], mapping: Dict):
    """Recursively map nodes."""
    if not isinstance(node, dict):
        return
        
    node_id = node.get("node_id")
    if node_id:
        mapping[node_id] = node
    
    for child in node.get("nodes", []):
        _map_node(child, mapping)


def get_all_text_from_tree(tree: Dict[str, Any]) -> str:
    """Extract all text content from a tree for full-text search."""
    texts = []
    
    def extract(node):
        if isinstance(node, dict):
            for key in ["text", "summary", "title"]:
                if key in node and node[key]:
                    texts.append(str(node[key]))
            for child in node.get("nodes", []):
                extract(child)
        elif isinstance(node, list):
            for item in node:
                extract(item)
    
    if "structure" in tree:
        extract(tree["structure"])
    else:
        extract(tree)
    
    return "\n\n".join(texts)
